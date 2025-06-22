import os
import json
import re
from collections import defaultdict
from pdfminer.high_level import extract_text
from docx import Document
from pptx import Presentation

CATEGORY_KEYWORDS = {
    "Web3 & AI": ["web3", "blockchain", "crypto", "ai", "artificial intelligence", "machine learning", "decentralized"],
    "Tech & EdTech": ["edtech", "education", "learning", "training", "tech", "software"],
    "E-commerce & Brands": ["e-commerce", "brand", "marketing", "shopping", "consumer", "retail", "digital shelf", "commerce"],
    "Exhibitions & B2B": ["exhibition", "expo", "conference", "b2b", "event"],
    "Entertainment & Social": ["entertainment", "social", "community", "influencer", "game", "content"],
}

def categorize(text):
    text_l = text.lower()
    for category, keywords in CATEGORY_KEYWORDS.items():
        for kw in keywords:
            if kw in text_l:
                return category
    return "Tech & EdTech"

def extract_pdf(path):
    try:
        text = extract_text(path)
    except Exception:
        return "", ""
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    title = lines[0] if lines else os.path.basename(path)
    body = " ".join(lines[1:])
    sentences = re.split(r"(?<=[.!?]) +", body)
    summary = " ".join(sentences[:3]).strip()
    return title, summary

def extract_docx(path):
    try:
        doc = Document(path)
        lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    except Exception:
        return "", ""
    title = lines[0] if lines else os.path.basename(path)
    body = " ".join(lines[1:])
    sentences = re.split(r"(?<=[.!?]) +", body)
    summary = " ".join(sentences[:3]).strip()
    return title, summary

def extract_pptx(path):
    try:
        pres = Presentation(path)
        texts = []
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text.strip())
    except Exception:
        return "", ""
    texts = [t for t in texts if t]
    title = texts[0] if texts else os.path.basename(path)
    body = " ".join(texts[1:])
    sentences = re.split(r"(?<=[.!?]) +", body)
    summary = " ".join(sentences[:3]).strip()
    return title, summary

EXTRACTORS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
}

projects = []
for fname in os.listdir('.'):
    ext = os.path.splitext(fname)[1].lower()
    if ext in EXTRACTORS:
        title, summary = EXTRACTORS[ext](fname)
        if not title:
            title = os.path.splitext(fname)[0]
        if not summary:
            summary = "Summary unavailable."
        category = categorize(title + " " + summary)
        projects.append({
            "title": title,
            "category": category,
            "summary": summary,
            "path": fname,
        })

os.makedirs('site', exist_ok=True)
with open('site/projects.json', 'w') as f:
    json.dump(projects, f, indent=2)

def slugify(text):
    return text.lower().replace(' & ', '-').replace(' ', '-')

# group projects by category
cats = defaultdict(list)
for p in projects:
    cats[p['category']].append(p)

# Base HTML templates
header = '''<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{title}</title>
<link href="https://fonts.googleapis.com/css2?family=Oxanium:wght@300;400;600&family=Space+Grotesk:wght@300;400;600&display=swap" rel="stylesheet">
<link rel="stylesheet" href="style.css">
</head>
<body>
<nav>
  <a href="index.html" class="logo">Srijan Mishra</a>
  <div class="links">
    <a href="index.html">Home</a>
    <a href="#about">About</a>
    {cat_links}
  </div>
</nav>
<div class="container">
'''

footer = '''</div>
<footer>
  <p>&copy; 2024 Srijan Mishra</p>
</footer>
<script src="script.js"></script>
</body>
</html>'''

style_css = '''body {margin:0;font-family:"Space Grotesk",sans-serif;background:linear-gradient(135deg,#141e30,#243b55);color:#fff;}
nav {display:flex;justify-content:space-between;align-items:center;padding:1rem 2rem;background:rgba(255,255,255,0.05);backdrop-filter:blur(10px);position:sticky;top:0;z-index:10;}
nav a {color:#fff;margin:0 0.5rem;text-decoration:none;position:relative;}
nav a:hover::after,nav a.active::after{content:"";position:absolute;left:0;bottom:-4px;height:2px;width:100%;background:#0ff;transition:0.3s;}
.container{padding:2rem;}
.hero{display:flex;flex-direction:column;align-items:center;justify-content:center;padding:4rem 2rem;text-align:center;background:radial-gradient(circle at top left,rgba(0,255,255,0.2),transparent),radial-gradient(circle at bottom right,rgba(255,0,255,0.2),transparent);border-radius:1rem;}
.card-grid{display:grid;grid-template-columns:repeat(auto-fit,minmax(250px,1fr));gap:1.5rem;margin-top:2rem;}
.card{background:rgba(255,255,255,0.05);backdrop-filter:blur(10px);padding:1.5rem;border-radius:1rem;position:relative;border:1px solid rgba(255,255,255,0.1);transition:transform 0.3s,border-color 0.3s;}
.card:hover{transform:translateY(-5px);border-color:#0ff;}
.btn{display:inline-block;margin-top:1rem;padding:0.5rem 1rem;border:1px solid #0ff;border-radius:0.5rem;text-decoration:none;color:#0ff;transition:background 0.3s;}
.btn:hover{background:#0ff;color:#000;}
footer{text-align:center;padding:2rem 0;color:#aaa;font-size:0.9rem;}
@media (max-width:600px){nav .links{display:flex;flex-direction:column;}}
'''

script_js = '''document.querySelectorAll('nav a').forEach(a=>{if(a.href===location.href)a.classList.add('active');});'''

# Write CSS and JS
with open('site/style.css','w') as f:
    f.write(style_css)
with open('site/script.js','w') as f:
    f.write(script_js)

# Build index.html
cat_links_html = ''.join([f'<a href="{slugify(c)}.html">{c}</a>' for c in cats])
index_html = header.format(title='Home', cat_links=cat_links_html)
index_html += '''<section class="hero">
<h1>Portfolio</h1>
<p>Exploring the intersection of Web3, AI, and modern technology.</p>
</section>
<section id="about">
<h2>About</h2>
<p>Welcome to my portfolio showcasing work across Web3, AI, tech, and more.</p>
</section>'''
for c, items in cats.items():
    index_html += f'<section><h2>{c}</h2><div class="card-grid">'
    for p in items:
        index_html += f'<div class="card"><h3>{p["title"]}</h3><p>{p["summary"]}</p><a class="btn" href="../{p["path"]}" target="_blank">View File</a></div>'
    index_html += '</div></section>'
index_html += footer
with open('site/index.html','w') as f:
    f.write(index_html)

# Build category pages
for c, items in cats.items():
    html = header.format(title=c, cat_links=cat_links_html)
    html += f'<h1>{c}</h1><div class="card-grid">'
    for p in items:
        html += f'<div class="card"><h3>{p["title"]}</h3><p>{p["summary"]}</p><a class="btn" href="../{p["path"]}" target="_blank">View File</a></div>'
    html += '</div>'
    html += footer
    with open(f'site/{slugify(c)}.html','w') as f:
        f.write(html)
